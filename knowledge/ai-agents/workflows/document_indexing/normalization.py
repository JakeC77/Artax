"""Document normalization: detect type, extract content, emit spans (Step 1)."""

from __future__ import annotations

import io
import logging
import uuid
from typing import Optional

from app.workflows.document_indexing.models import Locator, Span

logger = logging.getLogger(__name__)


def _detect_content_type(filename: str, content_type: Optional[str], raw: bytes) -> str:
    """Return normalized type: pdf | docx | pptx | xlsx | csv | txt | md."""
    ext = (filename or "").lower().split(".")[-1] if "." in (filename or "") else ""
    
    # Log detection inputs for debugging
    logger.debug(
        "Detecting content type: filename=%s extension=%s content_type=%s",
        filename or "(none)",
        ext or "(none)",
        content_type or "(none)",
    )
    
    # Check extension first (most reliable) before content_type
    if ext == "pptx":
        logger.debug("Detected PPTX from extension")
        return "pptx"
    if ext == "docx":
        logger.debug("Detected DOCX from extension")
        return "docx"
    if ext == "pdf":
        logger.debug("Detected PDF from extension")
        return "pdf"
    if ext in ("xlsx", "xls"):
        logger.debug("Detected XLSX from extension")
        return "xlsx"
    if ext == "csv":
        logger.debug("Detected CSV from extension")
        return "csv"
    if ext == "txt":
        logger.debug("Detected TXT from extension")
        return "txt"
    if ext in ("md", "markdown"):
        logger.debug("Detected MD from extension")
        return "md"
    
    # Then check content_type (less reliable, but useful when extension missing)
    if content_type:
        ct = content_type.lower()
        # Be more specific - check PowerPoint first to avoid conflicts
        if "powerpoint" in ct or "presentation" in ct or "vnd.ms-powerpoint" in ct or "vnd.openxmlformats-officedocument.presentationml" in ct:
            logger.debug("Detected PPTX from content_type: %s", content_type)
            return "pptx"
        if "pdf" in ct:
            logger.debug("Detected PDF from content_type: %s", content_type)
            return "pdf"
        if ("word" in ct or "document" in ct or "vnd.ms-word" in ct or "vnd.openxmlformats-officedocument.wordprocessingml" in ct) and ext != "pptx":
            # Only return docx if we're sure it's not pptx (already checked above)
            logger.debug("Detected DOCX from content_type: %s", content_type)
            return "docx"
        if "sheet" in ct or "excel" in ct or "vnd.ms-excel" in ct or "vnd.openxmlformats-officedocument.spreadsheetml" in ct:
            logger.debug("Detected XLSX from content_type: %s", content_type)
            return "xlsx"
        if "csv" in ct:
            logger.debug("Detected CSV from content_type: %s", content_type)
            return "csv"
        if "text" in ct or "plain" in ct:
            logger.debug("Detected TXT from content_type: %s", content_type)
            return "txt" if ext != "md" else "md"
    
    # Check file signature (magic bytes)
    if raw[:4] == b"%PDF":
        logger.debug("Detected PDF from file signature")
        return "pdf"
    
    # Default fallback
    logger.warning(
        "Could not detect content type from filename=%s extension=%s content_type=%s, defaulting to PDF",
        filename or "(none)",
        ext or "(none)",
        content_type or "(none)",
    )
    return "pdf"  # default


def _spans_from_pdf(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from PDF (PyPDF2)."""
    try:
        import PyPDF2
    except ImportError:
        error_msg = "PyPDF2 not installed; PDF extraction unavailable"
        logger.error(error_msg)
        raise ImportError(error_msg)
    
    logger.debug("Starting PDF extraction: doc_id=%s size=%d bytes", doc_id[:8] if doc_id else "", len(raw))
    spans: list[Span] = []
    
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(raw))
        logger.debug("PDF loaded: %d pages", len(reader.pages))
        
        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if not text.strip():
                logger.debug("Page %d: no text extracted", page_num)
                continue
            span_id = str(uuid.uuid4())
            locator = Locator(type="pdf", page=page_num)
            spans.append(
                Span(
                    span_id=span_id,
                    doc_id=doc_id,
                    tenant_id=tenant_id,
                    text=text.strip(),
                    locator=locator,
                )
            )
            logger.debug("Page %d: extracted %d characters", page_num, len(text.strip()))
        
        logger.info("PDF extraction completed: doc_id=%s pages=%d spans=%d", doc_id[:8] if doc_id else "", len(reader.pages), len(spans))
        return spans
        
    except Exception as e:
        error_msg = f"Failed to extract text from PDF: {e}"
        logger.error(error_msg, exc_info=True)
        raise


def _spans_from_docx(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from DOCX (python-docx)."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        error_msg = "python-docx not installed; DOCX extraction unavailable"
        logger.error(error_msg)
        raise ImportError(error_msg)
    
    logger.debug("Starting DOCX extraction: doc_id=%s size=%d bytes", doc_id[:8] if doc_id else "", len(raw))
    spans: list[Span] = []
    
    try:
        doc = DocxDocument(io.BytesIO(raw))
        logger.debug("DOCX document loaded: %d paragraphs", len(doc.paragraphs))
        
        for para_idx, para in enumerate(doc.paragraphs):
            text = (para.text or "").strip()
            if not text:
                continue
            span_id = str(uuid.uuid4())
            locator = Locator(type="docx", paragraph_index=para_idx)
            spans.append(
                Span(
                    span_id=span_id,
                    doc_id=doc_id,
                    tenant_id=tenant_id,
                    text=text,
                    locator=locator,
                )
            )
        
        logger.info("DOCX extraction completed: doc_id=%s paragraphs=%d spans=%d", doc_id[:8] if doc_id else "", len(doc.paragraphs), len(spans))
        return spans
        
    except Exception as e:
        error_msg = f"Failed to extract text from DOCX: {e}"
        logger.error(error_msg, exc_info=True)
        raise


def _spans_from_pptx(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from PPTX (python-pptx)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        error_msg = "python-pptx not installed; PPTX extraction unavailable"
        logger.error(error_msg)
        raise ImportError(error_msg)
    
    logger.debug("Starting PPTX extraction: doc_id=%s size=%d bytes", doc_id[:8] if doc_id else "", len(raw))
    spans: list[Span] = []
    
    try:
        prs = Presentation(io.BytesIO(raw))
        logger.debug("PPTX presentation loaded: %d slides", len(prs.slides))
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = (shape.text or "").strip()
                    if not text:
                        continue
                    span_id = str(uuid.uuid4())
                    shape_id = getattr(shape, "shape_id", None) or str(id(shape))
                    locator = Locator(type="pptx", slide=slide_idx, shape_id=shape_id)
                    spans.append(
                        Span(
                            span_id=span_id,
                            doc_id=doc_id,
                            tenant_id=tenant_id,
                            text=text,
                            locator=locator,
                        )
                    )
                    slide_text_count += 1
            logger.debug("Slide %d: extracted %d text spans", slide_idx, slide_text_count)
        
        logger.info("PPTX extraction completed: doc_id=%s slides=%d spans=%d", doc_id[:8] if doc_id else "", len(prs.slides), len(spans))
        return spans
        
    except Exception as e:
        error_msg = f"Failed to extract text from PPTX: {e}"
        logger.error(error_msg, exc_info=True)
        raise


def _spans_from_xlsx(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from XLSX (openpyxl)."""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl not installed; XLSX extraction unavailable")
        return []
    spans: list[Span] = []
    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            for col_idx, cell in enumerate(row, start=1):
                if cell is not None and str(cell).strip():
                    span_id = str(uuid.uuid4())
                    locator = Locator(type="xlsx", sheet=sheet_name, row=row_idx, column=col_idx)
                    spans.append(
                        Span(
                            span_id=span_id,
                            doc_id=doc_id,
                            tenant_id=tenant_id,
                            text=str(cell).strip(),
                            locator=locator,
                        )
                    )
    wb.close()
    return spans


def _spans_from_csv(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from CSV."""
    import csv
    spans: list[Span] = []
    try:
        text = raw.decode("utf-8-sig")
    except Exception:
        text = raw.decode("latin-1")
    reader = csv.reader(io.StringIO(text))
    for row_idx, row in enumerate(reader, start=1):
        for col_idx, cell in enumerate(row, start=1):
            if cell and cell.strip():
                span_id = str(uuid.uuid4())
                locator = Locator(type="csv", sheet="default", row=row_idx, column=col_idx)
                spans.append(
                    Span(
                        span_id=span_id,
                        doc_id=doc_id,
                        tenant_id=tenant_id,
                        text=cell.strip(),
                        locator=locator,
                    )
                )
    return spans


def _spans_from_txt(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from plain text file (split by paragraphs)."""
    spans: list[Span] = []
    try:
        text = raw.decode("utf-8-sig")
    except Exception:
        try:
            text = raw.decode("utf-8")
        except Exception:
            text = raw.decode("latin-1", errors="ignore")
    
    # Split by blank lines (paragraphs) or by single newlines if no blank lines
    paragraphs = text.split("\n\n")
    if len(paragraphs) == 1:
        # No blank lines, split by single newlines
        paragraphs = [p for p in text.split("\n") if p.strip()]
    
    for para_idx, para in enumerate(paragraphs, start=1):
        para_text = para.strip()
        if not para_text:
            continue
        span_id = str(uuid.uuid4())
        locator = Locator(type="txt", paragraph_index=para_idx)
        spans.append(
            Span(
                span_id=span_id,
                doc_id=doc_id,
                tenant_id=tenant_id,
                text=para_text,
                locator=locator,
            )
        )
    return spans


def _spans_from_md(raw: bytes, doc_id: str, tenant_id: str) -> list[Span]:
    """Extract spans from Markdown file (split by paragraphs, preserving structure)."""
    spans: list[Span] = []
    try:
        text = raw.decode("utf-8-sig")
    except Exception:
        try:
            text = raw.decode("utf-8")
        except Exception:
            text = raw.decode("latin-1", errors="ignore")
    
    # Split by blank lines (paragraphs) or by single newlines if no blank lines
    paragraphs = text.split("\n\n")
    if len(paragraphs) == 1:
        # No blank lines, split by single newlines
        paragraphs = [p for p in text.split("\n") if p.strip()]
    
    for para_idx, para in enumerate(paragraphs, start=1):
        para_text = para.strip()
        if not para_text:
            continue
        span_id = str(uuid.uuid4())
        locator = Locator(type="md", paragraph_index=para_idx)
        spans.append(
            Span(
                span_id=span_id,
                doc_id=doc_id,
                tenant_id=tenant_id,
                text=para_text,
                locator=locator,
            )
        )
    return spans


class UnsupportedFileTypeError(Exception):
    """Raised when a file type is not supported."""
    pass


class NormalizationError(Exception):
    """Raised when normalization fails."""
    pass


def normalize_to_spans(
    raw: bytes,
    doc_id: str,
    tenant_id: str,
    filename: Optional[str] = None,
    content_type: Optional[str] = None,
) -> list[Span]:
    """
    Detect file type, extract content, and return list of spans.
    Supports PDF, DOCX, PPTX, XLSX, CSV, TXT, MD.
    
    Raises:
        UnsupportedFileTypeError: If file type is not supported
        NormalizationError: If normalization fails
    """
    logger.info(
        "Starting normalization: doc_id=%s filename=%s content_type=%s size=%d bytes",
        doc_id[:8] if doc_id else "(none)",
        filename or "(none)",
        content_type or "(none)",
        len(raw),
    )
    
    doc_type = _detect_content_type(filename or "", content_type, raw)
    logger.info("Detected document type: %s for filename=%s", doc_type, filename or "(none)")
    
    try:
        if doc_type == "pdf":
            spans = _spans_from_pdf(raw, doc_id, tenant_id)
        elif doc_type == "docx":
            spans = _spans_from_docx(raw, doc_id, tenant_id)
        elif doc_type == "pptx":
            spans = _spans_from_pptx(raw, doc_id, tenant_id)
        elif doc_type == "xlsx":
            spans = _spans_from_xlsx(raw, doc_id, tenant_id)
        elif doc_type == "csv":
            spans = _spans_from_csv(raw, doc_id, tenant_id)
        elif doc_type == "txt":
            spans = _spans_from_txt(raw, doc_id, tenant_id)
        elif doc_type == "md":
            spans = _spans_from_md(raw, doc_id, tenant_id)
        else:
            # This should not happen if _detect_content_type works correctly
            error_msg = f"Unsupported document type: {doc_type} (filename: {filename or 'unknown'})"
            logger.error(error_msg)
            raise UnsupportedFileTypeError(error_msg)
        
        logger.info(
            "Normalization completed: doc_id=%s type=%s spans=%d",
            doc_id[:8] if doc_id else "(none)",
            doc_type,
            len(spans),
        )
        
        return spans
        
    except UnsupportedFileTypeError:
        # Re-raise unsupported file type errors
        raise
    except ImportError as e:
        # Missing required library
        error_msg = f"Required library not installed for {doc_type} files: {e}"
        logger.error(error_msg)
        raise NormalizationError(error_msg) from e
    except Exception as e:
        # Other normalization errors
        error_msg = f"Failed to normalize {doc_type} file (filename: {filename or 'unknown'}): {e}"
        logger.error(error_msg, exc_info=True)
        raise NormalizationError(error_msg) from e
